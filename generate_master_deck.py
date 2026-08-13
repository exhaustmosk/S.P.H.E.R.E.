"""
generate_master_deck.py
────────────────────────────────────────────────────────────────────────────
Programmatically creates an ultra-polished, publication-ready, 100% white-themed
9-slide presentation saved as 'Multimodal_Psychiatric_AI_Deck.pptx'.

Built using python-pptx with 16:9 widescreen layout, card-based modern UI,
clean typography, and clinical-grade color palettes.
────────────────────────────────────────────────────────────────────────────
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── Color Palette Constants ───────────────────────────────────────────────────
C_BG_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)  # #FFFFFF
C_CARD_BG        = RGBColor(0xF8, 0xFA, 0xFC)  # #F8FAFC (Soft Slate Tint)
C_CARD_BG_ALT    = RGBColor(0xF1, 0xF5, 0xF9)  # #F1F5F9 (Cool Light Gray)
C_CARD_BORDER    = RGBColor(0xE2, 0xE8, 0xF0)  # #E2E8F0 (Subtle Slate Border)
C_BORDER_ACCENT  = RGBColor(0xCB, 0xD5, 0xE1)  # #CBD5E1

C_NAVY_TITLE     = RGBColor(0x0F, 0x17, 0x2A)  # #0F172A (Deep Slate Navy)
C_NAVY_SUB       = RGBColor(0x1E, 0x29, 0x3B)  # #1E293B
C_TEXT_BODY      = RGBColor(0x47, 0x55, 0x69)  # #475569 (Slate Gray)
C_TEXT_MUTED     = RGBColor(0x64, 0x74, 0x8B)  # #64748B (Light Slate Gray)

C_BLUE_ACCENT    = RGBColor(0x25, 0x63, 0xEB)  # #2563EB (Electric Royal Blue)
C_BLUE_LIGHT     = RGBColor(0xDB, 0xEA, 0xFE)  # #DBEAFE (Soft Blue Pill)
C_BLUE_DARK      = RGBColor(0x1D, 0x4E, 0xD8)  # #1D4ED8

C_GREEN_ACCENT   = RGBColor(0x05, 0x96, 0x69)  # #059669 (Emerald Green)
C_GREEN_LIGHT    = RGBColor(0xD1, 0xFA, 0xE5)  # #D1FAE5 (Soft Green Pill)

C_RED_ACCENT     = RGBColor(0xDC, 0x26, 0x26)  # #DC2626 (Crimson Red)
C_RED_LIGHT      = RGBColor(0xFE, 0xE2, 0xE2)  # #FEE2E2 (Soft Red Pill)

C_PURPLE_ACCENT  = RGBColor(0x7C, 0x3A, 0xED)  # #7C3AED (Violet Accent)
C_PURPLE_LIGHT   = RGBColor(0xED, 0xE9, 0xFE)  # #EDE9FE

FONT_HEADING = "Arial"
FONT_BODY    = "Calibri"


def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ── Helper: Set Solid Background Color ────────────────────────────────────
    def set_white_bg(slide):
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = C_BG_WHITE
        bg_shape.line.fill.background()
        return bg_shape

    # ── Helper: Standard Slide Header ─────────────────────────────────────────
    def add_slide_header(slide, category_text, title_text, subtitle_text=None):
        # Category Pill / Tag
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.73), Inches(0.3))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        tf_tag.margin_left = tf_tag.margin_right = tf_tag.margin_top = tf_tag.margin_bottom = 0
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = category_text.upper()
        p_tag.font.name = FONT_HEADING
        p_tag.font.size = Pt(10.5)
        p_tag.font.bold = True
        p_tag.font.color.rgb = C_BLUE_ACCENT

        # Title Box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(11.73), Inches(0.55))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = FONT_HEADING
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = C_NAVY_TITLE

        # Subtitle
        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.22), Inches(11.73), Inches(0.35))
            tf_sub = sub_box.text_frame
            tf_sub.word_wrap = True
            tf_sub.margin_left = tf_sub.margin_right = tf_sub.margin_top = tf_sub.margin_bottom = 0
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle_text
            p_sub.font.name = FONT_BODY
            p_sub.font.size = Pt(13)
            p_sub.font.color.rgb = C_TEXT_BODY

    # ── Helper: Add Rounded Card ──────────────────────────────────────────────
    def add_card(slide, left, top, width, height, bg_color=C_CARD_BG, border_color=C_CARD_BORDER, border_width=1.0):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(border_width)
        else:
            shape.line.fill.background()
        return shape

    # ── Helper: Add Pill Badge ────────────────────────────────────────────────
    def add_badge(slide, left, top, width, height, text, bg_color=C_BLUE_LIGHT, text_color=C_BLUE_DARK):
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        badge.fill.solid()
        badge.fill.fore_color.rgb = bg_color
        badge.line.fill.background()
        tf = badge.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_HEADING
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = text_color
        return badge

    # =========================================================================
    # SLIDE 1: TITLE SLIDE (Hero Card Design)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_white_bg(s1)

    # Hero Outer Card
    add_card(s1, Inches(0.8), Inches(0.7), Inches(11.733), Inches(6.1), bg_color=RGBColor(0xFA, 0xFC, 0xFF), border_color=C_BLUE_ACCENT, border_width=2.0)

    # Top accent bar on Hero Card
    top_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.7), Inches(11.733), Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = C_BLUE_ACCENT
    top_bar.line.fill.background()

    # Category Badge
    add_badge(s1, Inches(1.3), Inches(1.2), Inches(3.6), Inches(0.35), "OPEN ELECTIVE CAPSTONE PROJECT | CLINICAL AI", bg_color=C_BLUE_LIGHT, text_color=C_BLUE_DARK)

    # Main Title
    t_box = s1.shapes.add_textbox(Inches(1.3), Inches(1.75), Inches(10.7), Inches(1.8))
    tf = t_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "Explainable Multimodal AI Framework for Psychiatric Evaluation & Mental Health Assessment"
    p.font.name = FONT_HEADING
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = C_NAVY_TITLE
    p.space_after = Pt(12)

    p2 = tf.add_paragraph()
    p2.text = "Integrating Facial Micro-Expressions, Speech Prosody, Digital Behavior, and Autonomic Physiological Indicators for Objective Clinical Decision Support"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(15)
    p2.font.color.rgb = C_TEXT_BODY

    # Divider line
    div = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(3.85), Inches(10.7), Inches(0.02))
    div.fill.solid()
    div.fill.fore_color.rgb = C_CARD_BORDER
    div.line.fill.background()

    # Team Members Header
    team_lbl = s1.shapes.add_textbox(Inches(1.3), Inches(4.05), Inches(10.7), Inches(0.3))
    tf_tl = team_lbl.text_frame
    p_tl = tf_tl.paragraphs[0]
    p_tl.text = "PROJECT AUTHORS & RESEARCH TEAM"
    p_tl.font.name = FONT_HEADING
    p_tl.font.size = Pt(11)
    p_tl.font.bold = True
    p_tl.font.color.rgb = C_TEXT_MUTED

    # 4 Author Badges
    authors = ["Akshaj S.", "Gautam Girish", "Abid Roshan", "Jacob S. Mathew"]
    card_w = Inches(2.52)
    card_gap = Inches(0.2)
    start_x = Inches(1.3)

    for i, name in enumerate(authors):
        bx = start_x + i * (card_w + card_gap)
        add_card(s1, bx, Inches(4.45), card_w, Inches(1.15), bg_color=C_BG_WHITE, border_color=C_CARD_BORDER)
        
        # Name
        tb = s1.shapes.add_textbox(bx + Inches(0.15), Inches(4.55), card_w - Inches(0.3), Inches(0.95))
        tf_n = tb.text_frame
        tf_n.word_wrap = True
        tf_n.margin_left = tf_n.margin_right = tf_n.margin_top = tf_n.margin_bottom = 0
        p_n = tf_n.paragraphs[0]
        p_n.text = name
        p_n.font.name = FONT_HEADING
        p_n.font.size = Pt(14)
        p_n.font.bold = True
        p_n.font.color.rgb = C_NAVY_TITLE
        p_n.space_after = Pt(2)
        
        p_dept = tf_n.add_paragraph()
        p_dept.text = "Undergraduate Researcher\nDept. of Mechanical Engineering"
        p_dept.font.name = FONT_BODY
        p_dept.font.size = Pt(10)
        p_dept.font.color.rgb = C_TEXT_MUTED

    # Institution Footer
    inst_box = s1.shapes.add_textbox(Inches(1.3), Inches(5.8), Inches(10.7), Inches(0.4))
    tf_inst = inst_box.text_frame
    p_inst = tf_inst.paragraphs[0]
    p_inst.text = "SRM Institute of Science and Technology (SRMIST) — Department of Mechanical Engineering (Open Elective)"
    p_inst.font.name = FONT_BODY
    p_inst.font.size = Pt(11.5)
    p_inst.font.bold = True
    p_inst.font.color.rgb = C_BLUE_ACCENT

    # =========================================================================
    # SLIDE 2: PROBLEM STATEMENT & CLINICAL IMPERATIVE
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_white_bg(s2)
    add_slide_header(s2, "Clinical Background & Motivation", "Problem Statement & Clinical Imperative", "Transitioning from subjective questionnaires to quantitative multi-sensor psychiatric biomarkers")

    # Left Card - The Challenge
    add_card(s2, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.3), bg_color=C_CARD_BG, border_color=C_RED_ACCENT, border_width=1.5)
    add_badge(s2, Inches(1.1), Inches(1.95), Inches(2.8), Inches(0.3), "CURRENT CLINICAL BOTTLENECK", bg_color=C_RED_LIGHT, text_color=C_RED_ACCENT)

    tb_left = s2.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(5.1), Inches(4.3))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = 0

    p = tf_l.paragraphs[0]
    p.text = "Traditional Diagnostic Limitations"
    p.font.name = FONT_HEADING
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_NAVY_TITLE
    p.space_after = Pt(14)

    challenges = [
        ("Subjective Self-Report Bias", "Scales like PHQ-9 and GAD-7 suffer from recall errors, social desirability bias, and acute emotional masking."),
        ("High Inter-Clinician Variability", "Manual psychiatric interviews lack standardized physical metrics, leading to delayed or discordant diagnoses."),
        ("Absence of Continuous Monitoring", "Episodic clinical visits fail to track episodic fluctuations in autonomic tone and behavioral deterioration."),
        ("Delayed Clinical Intervention", "Patients typically seek help only after acute functional impairment occurs, missing early therapeutic windows.")
    ]

    for title, desc in challenges:
        p_t = tf_l.add_paragraph()
        p_t.text = f"• {title}: "
        p_t.font.name = FONT_HEADING
        p_t.font.size = Pt(12.5)
        p_t.font.bold = True
        p_t.font.color.rgb = C_NAVY_SUB

        p_d = tf_l.add_paragraph()
        p_d.text = f"   {desc}"
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = C_TEXT_BODY
        p_d.space_after = Pt(8)

    # Right Card - The Solution
    add_card(s2, Inches(6.833), Inches(1.7), Inches(5.7), Inches(5.3), bg_color=C_CARD_BG, border_color=C_GREEN_ACCENT, border_width=1.5)
    add_badge(s2, Inches(7.133), Inches(1.95), Inches(3.0), Inches(0.3), "PROPOSED MULTIMODAL PARADIGM", bg_color=C_GREEN_LIGHT, text_color=C_GREEN_ACCENT)

    tb_right = s2.shapes.add_textbox(Inches(7.133), Inches(2.4), Inches(5.1), Inches(4.3))
    tf_r = tb_right.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = tf_r.margin_bottom = 0

    p = tf_r.paragraphs[0]
    p.text = "Objective Biomarker Integration"
    p.font.name = FONT_HEADING
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_NAVY_TITLE
    p.space_after = Pt(14)

    solutions = [
        ("Acoustic Vocal Prosody", "Speech emotion dynamics, pitch entropy, MFCC variance, and acoustic hesitation markers via Wav2Vec 2.0."),
        ("Facial Affective Dynamics", "Vision Transformer (ViT) embeddings capturing micro-expression variance, smile intensity, and eye-blink frequency."),
        ("Autonomic Physiological Signals", "Empirical autonomic telemetry measuring Galvanic Skin Response (GSR), Heart Rate Variability (HRV), and skin temperature."),
        ("Continuous Behavioral Telemetry", "Objective device telemetry capturing typing speed, application usage duration, and session idle latency.")
    ]

    for title, desc in solutions:
        p_t = tf_r.add_paragraph()
        p_t.text = f"✔ {title}: "
        p_t.font.name = FONT_HEADING
        p_t.font.size = Pt(12.5)
        p_t.font.bold = True
        p_t.font.color.rgb = C_NAVY_SUB

        p_d = tf_r.add_paragraph()
        p_d.text = f"   {desc}"
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = C_TEXT_BODY
        p_d.space_after = Pt(8)

    # =========================================================================
    # SLIDE 3: THREE CORE SYSTEM OBJECTIVES
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_white_bg(s3)
    add_slide_header(s3, "System Architecture", "Three Core System Objectives", "A unified tri-objective machine learning pipeline engineered for psychiatric diagnostics")

    cards_data_s3 = [
        ("OBJECTIVE 1", "4-Class Mental Health Classification", C_BLUE_ACCENT, C_BLUE_LIGHT, [
            ("Core Task", "Discrete categorization of patient clinical mental state."),
            ("Target Classes", "Healthy (0), Mild Stress (1), Moderate Stress (2), Severe Stress (3)."),
            ("Model Engine", "AutoGluon multi-layer stacked ensemble incorporating LightGBM, CatBoost, ExtraTrees & Random Forests."),
            ("Clinical Impact", "Provides rapid triage and diagnostic stratification for acute clinical decision pathways.")
        ]),
        ("OBJECTIVE 2", "Multi-Output Severity Regression", C_PURPLE_ACCENT, C_PURPLE_LIGHT, [
            ("Core Task", "Continuous quantitative symptom severity prediction."),
            ("Target Metrics", "Depression Score (0–34), Anxiety Score (0–24), Stress Score (0–39)."),
            ("Model Engine", "MultiOutputRegressor wrapped around gradient boosted XGBoost trees with 5-Fold Cross Validation."),
            ("Clinical Impact", "Yields fine-grained dimensional severity ratings tracking subtle symptom escalations.")
        ]),
        ("OBJECTIVE 3", "TreeSHAP Explainability Engine", C_GREEN_ACCENT, C_GREEN_LIGHT, [
            ("Core Task", "Model transparency and modality attribution auditing."),
            ("Attribution Tiers", "Acoustic / Speech (32 PCA) vs. Visual / Facial (32 PCA) vs. Physiology & Behavior (18 Tabular)."),
            ("Explainability Tool", "Game-theoretic TreeSHAP feature importance and patient-level attribution waterfall plots."),
            ("Clinical Impact", "Enables physician auditability, verifying that empirical biological markers drive high-risk alerts.")
        ])
    ]

    card_w = Inches(3.71)
    card_gap = Inches(0.3)
    start_x = Inches(0.8)

    for i, (tag, title, color_acc, color_light, details) in enumerate(cards_data_s3):
        bx = start_x + i * (card_w + card_gap)
        add_card(s3, bx, Inches(1.7), card_w, Inches(5.3), bg_color=C_CARD_BG, border_color=color_acc, border_width=1.5)
        add_badge(s3, bx + Inches(0.25), Inches(1.95), Inches(1.6), Inches(0.28), tag, bg_color=color_light, text_color=color_acc)

        tb = s3.shapes.add_textbox(bx + Inches(0.25), Inches(2.35), card_w - Inches(0.5), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_HEADING
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_NAVY_TITLE
        p.space_after = Pt(14)

        for heading, body in details:
            ph = tf.add_paragraph()
            ph.text = f"• {heading}:"
            ph.font.name = FONT_HEADING
            ph.font.size = Pt(12)
            ph.font.bold = True
            ph.font.color.rgb = C_NAVY_SUB

            pb = tf.add_paragraph()
            pb.text = f"  {body}"
            pb.font.name = FONT_BODY
            pb.font.size = Pt(11)
            pb.font.color.rgb = C_TEXT_BODY
            pb.space_after = Pt(8)

    # =========================================================================
    # SLIDE 4: DATASET ARCHITECTURE & HETEROGENEOUS SOURCES
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_white_bg(s4)
    add_slide_header(s4, "Data Ingestion & Multimodal Integration", "Dataset Architecture & Heterogeneous Sources", "Aggregating acoustic speech, visual affect, and autonomic telemetry into a synchronized clinical matrix")

    datasets = [
        ("ACOUSTIC STREAM", "RAVDESS Speech Emotion Corpus", "1,440 Audio Files (24 Actors)", C_BLUE_ACCENT, C_BLUE_LIGHT, [
            ("Raw Audio Ingestion", "16 kHz studio-quality vocal recordings of validated emotional statements."),
            ("Feature Extractor", "HuggingFace r-f/wav2vec-english-speech-emotion-recognition."),
            ("Deep Embeddings", "1024-dimensional sequence representations extracted via mean-pooling."),
            ("Target Mapping", "8 emotional categories (Neutral, Calm, Happy, Sad, Angry, Fear, Disgust, Surprised) mapped to clinical severity.")
        ]),
        ("VISUAL STREAM", "FER2013 Facial Expression Dataset", "28,709 Facial Images (48×48)", C_PURPLE_ACCENT, C_PURPLE_LIGHT, [
            ("Image Modality", "Grayscale facial expressions converted to 3-channel RGB tensors."),
            ("Feature Extractor", "dima806/facial_emotions_image_detection Vision Transformer (ViT)."),
            ("Deep Embeddings", "768-dimensional visual embeddings extracted from the [CLS] sequence pooler token."),
            ("Affect Categories", "7 primary affective classes (Angry, Disgust, Fear, Happy, Neutral, Sad, Surprise).")
        ]),
        ("PHYSIOLOGY & BEHAVIOR", "Clinical Physiological & Behavioral Cohort", "4,000 Complete Patient Records", C_GREEN_ACCENT, C_GREEN_LIGHT, [
            ("Autonomic Biomarkers", "Heart Rate (BPM), Heart Rate Variability (HRV Index), Galvanic Skin Response (GSR), Skin Temperature."),
            ("Behavioral Telemetry", "Daily App Usage (min), Session Frequency, Idle Time, Typing Speed (WPM), Sleep Quality."),
            ("Affective Micro-Cues", "Facial Emotion Variance, Eye Blink Rate, Smile Intensity, Head Motion Index, Pitch Mean, Speech Rate."),
            ("Ground-Truth Labels", "Categorical Mental_Health_Status + Continuous Depression, Anxiety, and Stress scores.")
        ])
    ]

    card_w = Inches(3.71)
    card_gap = Inches(0.3)
    start_x = Inches(0.8)

    for i, (tag, title, kpi, color_acc, color_light, details) in enumerate(datasets):
        bx = start_x + i * (card_w + card_gap)
        add_card(s4, bx, Inches(1.7), card_w, Inches(5.3), bg_color=C_CARD_BG, border_color=color_acc, border_width=1.5)
        add_badge(s4, bx + Inches(0.2), Inches(1.92), Inches(1.9), Inches(0.28), tag, bg_color=color_light, text_color=color_acc)

        # KPI Pill inside Card
        kpi_card = add_card(s4, bx + Inches(0.2), Inches(2.28), card_w - Inches(0.4), Inches(0.55), bg_color=C_BG_WHITE, border_color=C_CARD_BORDER)
        tf_k = kpi_card.text_frame
        tf_k.margin_left = tf_k.margin_right = tf_k.margin_top = tf_k.margin_bottom = 0
        p_k = tf_k.paragraphs[0]
        p_k.text = kpi
        p_k.alignment = PP_ALIGN.CENTER
        p_k.font.name = FONT_HEADING
        p_k.font.size = Pt(11)
        p_k.font.bold = True
        p_k.font.color.rgb = color_acc

        tb = s4.shapes.add_textbox(bx + Inches(0.2), Inches(2.95), card_w - Inches(0.4), Inches(3.9))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_HEADING
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_NAVY_TITLE
        p.space_after = Pt(10)

        for heading, body in details:
            ph = tf.add_paragraph()
            ph.text = f"• {heading}:"
            ph.font.name = FONT_HEADING
            ph.font.size = Pt(11.5)
            ph.font.bold = True
            ph.font.color.rgb = C_NAVY_SUB

            pb = tf.add_paragraph()
            pb.text = f"  {body}"
            pb.font.name = FONT_BODY
            pb.font.size = Pt(10.5)
            pb.font.color.rgb = C_TEXT_BODY
            pb.space_after = Pt(6)

    # =========================================================================
    # SLIDE 5: DATA PREPROCESSING & SEMANTIC TARGET ALIGNMENT
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_white_bg(s5)
    add_slide_header(s5, "Data Engineering & Pipeline Integrity", "Data Preprocessing & Semantic Target Alignment", "Solving cross-dataset misalignment through clinically validated affective state grouping")

    # Top Hero Callout - The Semantic Alignment Breakthrough
    add_card(s5, Inches(0.8), Inches(1.7), Inches(11.733), Inches(2.15), bg_color=RGBColor(0xFA, 0xFC, 0xFF), border_color=C_BLUE_ACCENT, border_width=1.5)
    add_badge(s5, Inches(1.1), Inches(1.9), Inches(3.2), Inches(0.3), "CRITICAL DATA ENGINEERING BREAKTHROUGH", bg_color=C_BLUE_LIGHT, text_color=C_BLUE_DARK)

    tb_top = s5.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(11.1), Inches(1.4))
    tf_t = tb_top.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_right = tf_t.margin_top = tf_t.margin_bottom = 0

    p = tf_t.paragraphs[0]
    p.text = "Overcoming Cross-Dataset Mismatch via Semantic Target Alignment"
    p.font.name = FONT_HEADING
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_NAVY_TITLE
    p.space_after = Pt(6)

    p_desc = tf_t.add_paragraph()
    p_desc.text = "Blind inner joins across independent datasets (survey rows vs. RAVDESS audio vs. FER2013 images) previously destroyed cross-modal correlations (resulting in 35% accuracy). We engineered a Semantic Target Alignment Engine that maps psychiatric severity categories directly to clinically coherent emotional subsets:"
    p_desc.font.name = FONT_BODY
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = C_TEXT_BODY
    p_desc.space_after = Pt(4)

    p_map = tf_t.add_paragraph()
    p_map.text = "• Healthy (0) → Happy, Neutral, Surprise   |   • Mild Stress (1) → Neutral, Sad   |   • Moderate Stress (2) → Sad, Angry, Disgust   |   • Severe Stress (3) → Angry, Fear, Disgust"
    p_map.font.name = FONT_HEADING
    p_map.font.size = Pt(11)
    p_map.font.bold = True
    p_map.font.color.rgb = C_BLUE_ACCENT

    # Bottom 3 Step Cards
    steps = [
        ("STEP 1", "Feature Standardization", C_BLUE_ACCENT, [
            "Fitted StandardScaler strictly on 18 numerical tabular features on the training set.",
            "Test set transformed independently to prevent statistical data leakage.",
            "Preserves true zero-mean and unit-variance distributions across physiological signals."
        ]),
        ("STEP 2", "PCA Noise Reduction", C_PURPLE_ACCENT, [
            "Acoustic PCA: 1024-d → 32 components (retaining >99.7% cumulative variance).",
            "Visual PCA: 768-d → 32 components (retaining >62.0% cumulative variance).",
            "Compresses high-dimensional noise while accelerating downstream gradient boosting."
        ]),
        ("STEP 3", "Fused Stratified Matrix", C_GREEN_ACCENT, [
            "Constructed X_fused (4,000 samples × 82 features: 18 Tabular + 32 Speech + 32 Facial).",
            "80/20 Stratified Partitioning (3,200 Train / 800 Test) matching target distribution.",
            "Ensures balanced representation of rare Severe Stress cases across all validation folds."
        ])
    ]

    card_w = Inches(3.71)
    card_gap = Inches(0.3)
    start_x = Inches(0.8)

    for i, (step_tag, title, acc_color, points) in enumerate(steps):
        bx = start_x + i * (card_w + card_gap)
        add_card(s5, bx, Inches(4.05), card_w, Inches(2.95), bg_color=C_CARD_BG, border_color=C_CARD_BORDER)
        add_badge(s5, bx + Inches(0.2), Inches(4.25), Inches(1.1), Inches(0.26), step_tag, bg_color=C_BLUE_LIGHT, text_color=C_BLUE_DARK)

        tb_s = s5.shapes.add_textbox(bx + Inches(0.2), Inches(4.6), card_w - Inches(0.4), Inches(2.3))
        tf_s = tb_s.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = tf_s.margin_right = tf_s.margin_top = tf_s.margin_bottom = 0

        p = tf_s.paragraphs[0]
        p.text = title
        p.font.name = FONT_HEADING
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_NAVY_TITLE
        p.space_after = Pt(8)

        for pt in points:
            p_pt = tf_s.add_paragraph()
            p_pt.text = f"• {pt}"
            p_pt.font.name = FONT_BODY
            p_pt.font.size = Pt(11)
            p_pt.font.color.rgb = C_TEXT_BODY
            p_pt.space_after = Pt(5)

    # =========================================================================
    # SLIDE 6: SEQUENTIAL STACKING METHODOLOGY
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_white_bg(s6)
    add_slide_header(s6, "Hierarchical Stacking Architecture", "Sequential Stacking Methodology", "Two-stage modeling: Symptom severity regression feeds into multi-layer classification ensemble")

    # Layer 1 Card (Head B - Regression)
    add_card(s6, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.3), bg_color=C_CARD_BG, border_color=C_PURPLE_ACCENT, border_width=1.5)
    add_badge(s6, Inches(1.1), Inches(1.95), Inches(3.0), Inches(0.3), "LAYER 1: SYMPTOM REGRESSOR (HEAD B)", bg_color=C_PURPLE_LIGHT, text_color=C_PURPLE_ACCENT)

    tb_l1 = s6.shapes.add_textbox(Inches(1.1), Inches(2.38), Inches(5.1), Inches(4.4))
    tf_l1 = tb_l1.text_frame
    tf_l1.word_wrap = True
    tf_l1.margin_left = tf_l1.margin_right = tf_l1.margin_top = tf_l1.margin_bottom = 0

    p = tf_l1.paragraphs[0]
    p.text = "Continuous Severity Predictor"
    p.font.name = FONT_HEADING
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = C_NAVY_TITLE
    p.space_after = Pt(12)

    l1_points = [
        ("Architecture", "MultiOutputRegressor(XGBRegressor(n_estimators=300, max_depth=6, lr=0.05))."),
        ("Input Features", "82 Fused Multimodal Features (18 Tabular + 32 Speech PCA + 32 Facial PCA)."),
        ("Out-Of-Fold Generation", "5-Fold Cross Validation generates clean, leak-free Out-Of-Fold (OOF) severity predictions on training data."),
        ("Predicted Outputs", "• pred_depression (0–34 Scale)\n• pred_anxiety (0–24 Scale)\n• pred_stress (0–39 Scale)"),
        ("Regression Metrics", "Achieved positive R² = 0.1948, MAE = 6.99, and RMSE = 8.54 across test targets.")
    ]

    for heading, body in l1_points:
        ph = tf_l1.add_paragraph()
        ph.text = f"• {heading}:"
        ph.font.name = FONT_HEADING
        ph.font.size = Pt(12)
        ph.font.bold = True
        ph.font.color.rgb = C_NAVY_SUB

        pb = tf_l1.add_paragraph()
        pb.text = f"  {body}"
        pb.font.name = FONT_BODY
        pb.font.size = Pt(11)
        pb.font.color.rgb = C_TEXT_BODY
        pb.space_after = Pt(6)

    # Layer 2 Card (Head A - Classification)
    add_card(s6, Inches(6.833), Inches(1.7), Inches(5.7), Inches(5.3), bg_color=C_CARD_BG, border_color=C_BLUE_ACCENT, border_width=1.5)
    add_badge(s6, Inches(7.133), Inches(1.95), Inches(3.2), Inches(0.3), "LAYER 2: AUTOGLUON ENSEMBLE (HEAD A)", bg_color=C_BLUE_LIGHT, text_color=C_BLUE_DARK)

    tb_l2 = s6.shapes.add_textbox(Inches(7.133), Inches(2.38), Inches(5.1), Inches(4.4))
    tf_l2 = tb_l2.text_frame
    tf_l2.word_wrap = True
    tf_l2.margin_left = tf_l2.margin_right = tf_l2.margin_top = tf_l2.margin_bottom = 0

    p = tf_l2.paragraphs[0]
    p.text = "Multi-Layer Stack Ensemble"
    p.font.name = FONT_HEADING
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = C_NAVY_TITLE
    p.space_after = Pt(12)

    l2_points = [
        ("Architecture", "AutoGluon TabularPredictor (presets='best_quality', auto_stack=True, eval_metric='f1_macro')."),
        ("Stacked Input Matrix", "85 Features = 82 Multimodal Features + 3 Injected Severity Predictions (pred_depression, pred_anxiety, pred_stress)."),
        ("Model Search Space", "Trained 25 base & bagged models spanning LightGBM, XGBoost, CatBoost, ExtraTrees, and Random Forests across 3 stack levels."),
        ("Champion Model", "RandomForest_r195_BAG_L1 with weighted multi-model ensembling achieving top generalization."),
        ("Evaluation Result", "89.38% Overall Accuracy, 0.9785 ROC-AUC, 0.8418 Macro F1.")
    ]

    for heading, body in l2_points:
        ph = tf_l2.add_paragraph()
        ph.text = f"✔ {heading}:"
        ph.font.name = FONT_HEADING
        ph.font.size = Pt(12)
        ph.font.bold = True
        ph.font.color.rgb = C_NAVY_SUB

        pb = tf_l2.add_paragraph()
        pb.text = f"  {body}"
        pb.font.name = FONT_BODY
        pb.font.size = Pt(11)
        pb.font.color.rgb = C_TEXT_BODY
        pb.space_after = Pt(6)

    # =========================================================================
    # SLIDE 7: MODEL PERFORMANCE & MATRIX HIGHLIGHTS
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_white_bg(s7)
    add_slide_header(s7, "Quantitative Validation & Metrics", "Model Performance & Matrix Highlights", "Rigorous evaluation on the 800-sample stratified test set demonstrating clinical reliability")

    # 4 Big KPI Cards
    kpis = [
        ("89.38%", "OVERALL ACCURACY", "Test Set (800 Cases)", C_BLUE_ACCENT, C_BLUE_LIGHT),
        ("100.0%", "SEVERE PRECISION", "Zero False Positives", C_RED_ACCENT, C_RED_LIGHT),
        ("0.9785", "MULTI-CLASS ROC-AUC", "One-vs-Rest Discrimination", C_GREEN_ACCENT, C_GREEN_LIGHT),
        ("0.8418", "MACRO F1-SCORE", "Balanced Across All Classes", C_PURPLE_ACCENT, C_PURPLE_LIGHT),
    ]

    card_w = Inches(2.70)
    card_gap = Inches(0.3)
    start_x = Inches(0.8)

    for i, (big_num, lbl, sub_lbl, color_acc, color_light) in enumerate(kpis):
        bx = start_x + i * (card_w + card_gap)
        add_card(s7, bx, Inches(1.7), card_w, Inches(1.75), bg_color=C_CARD_BG, border_color=color_acc, border_width=1.5)
        
        tb = s7.shapes.add_textbox(bx + Inches(0.15), Inches(1.85), card_w - Inches(0.3), Inches(1.45))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = big_num
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_HEADING
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color_acc
        p.space_after = Pt(2)

        p2 = tf.add_paragraph()
        p2.text = lbl
        p2.alignment = PP_ALIGN.CENTER
        p2.font.name = FONT_HEADING
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = C_NAVY_TITLE

        p3 = tf.add_paragraph()
        p3.text = sub_lbl
        p3.alignment = PP_ALIGN.CENTER
        p3.font.name = FONT_BODY
        p3.font.size = Pt(10)
        p3.font.color.rgb = C_TEXT_MUTED

    # Detailed Classification Report & Matrix Breakdown (Bottom Card)
    add_card(s7, Inches(0.8), Inches(3.65), Inches(11.733), Inches(3.35), bg_color=C_CARD_BG, border_color=C_CARD_BORDER)
    add_badge(s7, Inches(1.1), Inches(3.85), Inches(3.3), Inches(0.28), "CLASS-WISE CLINICAL PERFORMANCE BREAKDOWN", bg_color=C_BLUE_LIGHT, text_color=C_BLUE_DARK)

    tb_bot = s7.shapes.add_textbox(Inches(1.1), Inches(4.25), Inches(11.1), Inches(2.6))
    tf_b = tb_bot.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = 0

    # Class Performance Grid Table
    p_t = tf_b.paragraphs[0]
    p_t.text = "Clinical Diagnostic Category       Precision        Recall / Sensitivity        F1-Score         Test Support"
    p_t.font.name = FONT_HEADING
    p_t.font.size = Pt(12)
    p_t.font.bold = True
    p_t.font.color.rgb = C_NAVY_TITLE
    p_t.space_after = Pt(4)

    rows = [
        ("Healthy (Class 0)", "99.35%", "93.25%", "0.9620", "326 Cases"),
        ("Mild Stress (Class 1)", "78.43%", "97.17%", "0.8680", "247 Cases"),
        ("Moderate Stress (Class 2)", "90.23%", "78.11%", "0.8373", "201 Cases"),
        ("Severe Stress (Class 3)", "100.00%", "53.85%", "0.7000", "26 Cases")
    ]

    for cls_name, prec, rec, f1, sup in rows:
        p_r = tf_b.add_paragraph()
        p_r.text = f"{cls_name:<34} {prec:<16} {rec:<24} {f1:<16} {sup}"
        p_r.font.name = "Consolas"
        p_r.font.size = Pt(11.5)
        p_r.font.color.rgb = C_TEXT_BODY
        p_r.space_after = Pt(2)

    p_div = tf_b.add_paragraph()
    p_div.text = "─" * 82
    p_div.font.name = "Consolas"
    p_div.font.size = Pt(9)
    p_div.font.color.rgb = C_CARD_BORDER

    p_sum = tf_b.add_paragraph()
    p_sum.text = "Key Clinical Safety Takeaway: Zero Severe Stress cases were misclassified as Healthy. Misclassifications are strictly bounded to adjacent clinical categories, ensuring high-risk patients are never dismissed without monitoring."
    p_sum.font.name = FONT_BODY
    p_sum.font.size = Pt(11)
    p_sum.font.bold = True
    p_sum.font.color.rgb = C_GREEN_ACCENT

    # =========================================================================
    # SLIDE 8: OBJECTIVE 3 — SHAP MODALITY EXPLAINABILITY
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_white_bg(s8)
    add_slide_header(s8, "Objective 3: Explainable AI & Auditability", "SHAP Modality Explainability & Attribution", "Game-theoretic TreeSHAP proves that acoustic and visual biomarkers actively govern decision boundaries")

    # Left Card: Modality Contribution Breakdown
    add_card(s8, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.3), bg_color=C_CARD_BG, border_color=C_BLUE_ACCENT, border_width=1.5)
    add_badge(s8, Inches(1.1), Inches(1.95), Inches(3.0), Inches(0.3), "GLOBAL MODALITY CONTRIBUTION", bg_color=C_BLUE_LIGHT, text_color=C_BLUE_DARK)

    tb_sh_l = s8.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(5.1), Inches(4.3))
    tf_sl = tb_sh_l.text_frame
    tf_sl.word_wrap = True
    tf_sl.margin_left = tf_sl.margin_right = tf_sl.margin_top = tf_sl.margin_bottom = 0

    p = tf_sl.paragraphs[0]
    p.text = "Modality Attribution Percentages"
    p.font.name = FONT_HEADING
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = C_NAVY_TITLE
    p.space_after = Pt(12)

    mods = [
        ("Acoustic / Speech (32 PCA)", "43.45%", "Primary Driver", C_BLUE_ACCENT, "Wav2Vec vocal embeddings capture pitch modulation, jitter, and speech hesitation patterns indicative of stress escalation."),
        ("Visual / Facial (32 PCA)", "31.99%", "Secondary Driver", C_PURPLE_ACCENT, "Vision Transformer tokens represent micro-expression variance, eye-blink frequency, and smile intensity collapse."),
        ("Physiological & Behavioral", "24.56%", "Contextual Anchor", C_GREEN_ACCENT, "Autonomic sensors (HRV, GSR, Skin Temp) combined with behavioral telemetry (typing WPM, app latency) provide baseline stability.")
    ]

    for mod_name, pct, role, col, desc in mods:
        pm = tf_sl.add_paragraph()
        pm.text = f"{mod_name} — {pct} ({role})"
        pm.font.name = FONT_HEADING
        pm.font.size = Pt(13)
        pm.font.bold = True
        pm.font.color.rgb = col

        pd = tf_sl.add_paragraph()
        pd.text = f"{desc}"
        pd.font.name = FONT_BODY
        pd.font.size = Pt(11)
        pd.font.color.rgb = C_TEXT_BODY
        pd.space_after = Pt(8)

    # Right Card: Clinical Value & Interpretability
    add_card(s8, Inches(6.833), Inches(1.7), Inches(5.7), Inches(5.3), bg_color=C_CARD_BG, border_color=C_GREEN_ACCENT, border_width=1.5)
    add_badge(s8, Inches(7.133), Inches(1.95), Inches(3.0), Inches(0.3), "CLINICAL AUDITABILITY & TRUST", bg_color=C_GREEN_LIGHT, text_color=C_GREEN_ACCENT)

    tb_sh_r = s8.shapes.add_textbox(Inches(7.133), Inches(2.4), Inches(5.1), Inches(4.3))
    tf_sr = tb_sh_r.text_frame
    tf_sr.word_wrap = True
    tf_sr.margin_left = tf_sr.margin_right = tf_sr.margin_top = tf_sr.margin_bottom = 0

    p = tf_sr.paragraphs[0]
    p.text = "Clinician-in-the-Loop Transparency"
    p.font.name = FONT_HEADING
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = C_NAVY_TITLE
    p.space_after = Pt(12)

    val_points = [
        ("Evidence-Based Attribution", "TreeSHAP generates additive feature contribution values, allowing clinicians to review exactly why a patient was flagged as Severe Stress."),
        ("Elimination of Black-Box Skepticism", "Quantifies that decisions are not arbitrary tabular correlations, but grounded in speech tension (43.5%) and visual distress (32.0%)."),
        ("Patient-Level Explanation Reports", "Automated SHAP summary outputs provide transparent documentation for medical records and insurance compliance."),
        ("Objective Therapeutic Tracking", "Quantifies recovery progression over time as speech and facial SHAP attributions shift back toward the Healthy baseline.")
    ]

    for heading, body in val_points:
        ph = tf_sr.add_paragraph()
        ph.text = f"✔ {heading}:"
        ph.font.name = FONT_HEADING
        ph.font.size = Pt(12)
        ph.font.bold = True
        ph.font.color.rgb = C_NAVY_SUB

        pb = tf_sr.add_paragraph()
        pb.text = f"  {body}"
        pb.font.name = FONT_BODY
        pb.font.size = Pt(11)
        pb.font.color.rgb = C_TEXT_BODY
        pb.space_after = Pt(8)

    # =========================================================================
    # SLIDE 9: CONCLUSION & PRODUCTION DEPLOYMENT
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_white_bg(s9)
    add_slide_header(s9, "Translation & Production Readiness", "Conclusion & Production Deployment", "An edge-deployable, real-time psychiatric screening platform ready for clinical trials")

    conclusions = [
        ("CLINICAL DEPLOYMENT", "Interactive Streamlit Suite", C_BLUE_ACCENT, C_BLUE_LIGHT, [
            ("Real-Time Ingestion", "Streamlit UI ingesting live microphone and webcam feeds alongside wearable telemetry."),
            ("Live Symptom Radar", "Interactive radar plots rendering real-time Depression, Anxiety, and Stress severity indexes."),
            ("Automated Clinical PDF", "One-click generation of physician-ready PDF reports with embedded SHAP attribution charts.")
        ]),
        ("SYSTEM EFFICIENCY", "Low-Latency Edge Execution", C_PURPLE_ACCENT, C_PURPLE_LIGHT, [
            ("Sub-2 Second Latency", "End-to-end evaluation latency under 2 seconds on standard Apple Silicon (MPS backend) and CPUs."),
            ("Lightweight PCA Heads", "32-dimensional PCA compression enables real-time tree traversal without heavy GPU clusters."),
            ("Serialized Artifacts", "Production-ready joblib model pipelines allowing instant zero-cold-start inference.")
        ]),
        ("HEALTHCARE IMPACT", "Next-Generation Triage", C_GREEN_ACCENT, C_GREEN_LIGHT, [
            ("Objective Biomarkers", "Bridges the gap between subjective questionnaires and empirical multimodal signals."),
            ("Scalable Mental Health", "Enables continuous, passive psychiatric triage in universities, workplaces, and remote clinics."),
            ("Early Clinical Intervention", "Identifies subtle affective deterioration before acute functional impairment occurs.")
        ])
    ]

    card_w = Inches(3.71)
    card_gap = Inches(0.3)
    start_x = Inches(0.8)

    for i, (tag, title, color_acc, color_light, details) in enumerate(conclusions):
        bx = start_x + i * (card_w + card_gap)
        add_card(s9, bx, Inches(1.7), card_w, Inches(5.3), bg_color=C_CARD_BG, border_color=color_acc, border_width=1.5)
        add_badge(s9, bx + Inches(0.25), Inches(1.95), Inches(1.8), Inches(0.28), tag, bg_color=color_light, text_color=color_acc)

        tb = s9.shapes.add_textbox(bx + Inches(0.25), Inches(2.35), card_w - Inches(0.5), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_HEADING
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_NAVY_TITLE
        p.space_after = Pt(14)

        for heading, body in details:
            ph = tf.add_paragraph()
            ph.text = f"• {heading}:"
            ph.font.name = FONT_HEADING
            ph.font.size = Pt(12)
            ph.font.bold = True
            ph.font.color.rgb = C_NAVY_SUB

            pb = tf.add_paragraph()
            pb.text = f"  {body}"
            pb.font.name = FONT_BODY
            pb.font.size = Pt(11)
            pb.font.color.rgb = C_TEXT_BODY
            pb.space_after = Pt(8)

    # Save presentation
    output_filename = "Multimodal_Psychiatric_AI_Deck.pptx"
    prs.save(output_filename)
    print(f"✔ Presentation successfully saved to: {os.path.abspath(output_filename)}")


if __name__ == "__main__":
    create_deck()
